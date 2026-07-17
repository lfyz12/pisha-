"""Tests for ai_assistant.services.parsing."""

import io
import tempfile
from pathlib import Path

from django.test import SimpleTestCase

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from ai_assistant.services.parsing import extract_text


def build_pdf_bytes(text: str) -> bytes:
    """Build a minimal valid single-page PDF with extractable text."""
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = io.BytesIO()
    output.write(b"%PDF-1.4\n")
    offsets = []
    for number, obj in enumerate(objects, start=1):
        offsets.append(output.tell())
        output.write(f"{number} 0 obj\n".encode())
        output.write(obj)
        output.write(b"\nendobj\n")
    xref_pos = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n".encode())
    output.write(b"0000000000 65535 f \n")
    for offset in offsets:
        output.write(f"{offset:010d} 00000 n \n".encode())
    output.write(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF".encode()
    )
    return output.getvalue()


def build_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("First paragraph of the document.")
    document.add_paragraph("Second paragraph of the document.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Cell alpha"
    table.cell(0, 1).text = "Cell beta"
    table.cell(1, 1).text = "Cell omega"
    document.save(str(path))


def build_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly report"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    box.text_frame.text = "Revenue grew steadily this year."
    slide.notes_slide.notes_text_frame.text = "Mention the new scholarship program."
    presentation.save(str(path))


class ExtractTextTests(SimpleTestCase):
    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self._tmp.cleanup)
        self.dir = Path(self._tmp.name)

    def test_markdown_file_is_read_as_utf8(self):
        path = self.dir / "notes.md"
        path.write_text("# Консультант УПИШ\n\nТекст консультации.", encoding="utf-8")
        text = extract_text(path)
        self.assertIn("Консультант УПИШ", text)
        self.assertIn("Текст консультации.", text)

    def test_txt_file_with_invalid_utf8_bytes_is_read_with_replacement(self):
        path = self.dir / "broken.txt"
        path.write_bytes(b"valid prefix \xff\xfe invalid bytes")
        text = extract_text(path)
        self.assertIn("valid prefix", text)
        self.assertIn("�", text)

    def test_excessive_blank_lines_are_collapsed(self):
        path = self.dir / "blank.txt"
        path.write_text("line one\n\n\n\n\nline two", encoding="utf-8")
        text = extract_text(path)
        self.assertNotIn("\n\n\n", text)
        self.assertIn("line one\n\nline two", text)

    def test_docx_paragraphs_and_table_cells_are_extracted(self):
        path = self.dir / "document.docx"
        build_docx(path)
        text = extract_text(path)
        self.assertIn("First paragraph of the document.", text)
        self.assertIn("Second paragraph of the document.", text)
        self.assertIn("Cell alpha", text)
        self.assertIn("Cell omega", text)

    def test_pdf_text_is_extracted(self):
        path = self.dir / "document.pdf"
        path.write_bytes(build_pdf_bytes("Hello UPISH consultant"))
        self.assertIn("Hello UPISH consultant", extract_text(path))

    def test_pptx_shapes_and_speaker_notes_are_extracted(self):
        path = self.dir / "slides.pptx"
        build_pptx(path)
        text = extract_text(path)
        self.assertIn("Quarterly report", text)
        self.assertIn("Revenue grew steadily this year.", text)
        self.assertIn("Mention the new scholarship program.", text)

    def test_unsupported_extension_raises_value_error(self):
        path = self.dir / "archive.xyz"
        path.write_text("content", encoding="utf-8")
        with self.assertRaisesMessage(ValueError, "unsupported file extension"):
            extract_text(path)

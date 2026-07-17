"""Text extraction from uploaded knowledge-base documents."""

from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

#: Hard cap on the number of PDF pages processed for a single document.
MAX_PDF_PAGES = 200

_BLANK_LINES_RE = re.compile(r"\n{3,}")


class DocumentParsingError(ValueError):
    """Raised when a document cannot be parsed (corrupt or unreadable file).

    Subclasses ValueError so existing ``except ValueError`` callers keep
    working; unsupported extensions also raise it (see extract_text).
    """


def _normalize_blank_lines(text: str) -> str:
    """Collapse runs of 3+ newlines into a single blank line."""
    return _BLANK_LINES_RE.sub("\n\n", text).strip()


def _extract_plain(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_docx(path: Path) -> str:
    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages[:MAX_PDF_PAGES])


def _extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(parts)


_EXTRACTORS = {
    ".md": _extract_plain,
    ".txt": _extract_plain,
    ".docx": _extract_docx,
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
}


def extract_text(path) -> str:
    """Extract plain text from a document, dispatching on its file extension.

    Supported formats: .md, .txt (read as UTF-8 with replacement of invalid
    bytes), .docx (paragraphs and table cells), .pdf (page by page, capped at
    MAX_PDF_PAGES pages), .pptx (text frames of all slide shapes plus speaker
    notes). Runs of 3+ newlines in the result are collapsed to 2.

    Raises:
        DocumentParsingError: if the extension is not supported, or the file
            is corrupt/unreadable for its format (chained from the underlying
            library error). It subclasses ValueError, so ``except ValueError``
            catches both failure modes.
    """
    suffix = Path(path).suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise DocumentParsingError("unsupported file extension")
    try:
        text = extractor(Path(path))
    except Exception as exc:
        raise DocumentParsingError(
            f"failed to parse {suffix} document: {exc}"
        ) from exc
    return _normalize_blank_lines(text)
